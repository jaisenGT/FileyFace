# fileName, fileContent, webAddress

from pathlib import Path
import PyPDF2
import textrazor
import openai
import time
import os
from docx import Document
from pptx import Presentation
from ChromeHistory import get_web_history
from dotenv import load_dotenv

load_dotenv()

textrazor.api_key = "d6c3be0ed8bfbf0dcb235a80476a8d7dd009e91a316fd21b9f939563"
client = textrazor.TextRazor(extractors=["entities", "topics", "words", "phrases"])

openai.api_key = os.getenv('OPENAI_API_KEY')

def fileName(path):
    filePath = Path(path)
    fileName = filePath.name
    return fileName

def fileExtension(path):
    filePath = Path(path)
    fileExtension = filePath.suffix
    return fileExtension

# def open_file_with_retry(path, mode='rb', retries=5, delay=2):
#     for attempt in range(retries):
#         try:
#             return open(path, mode)
#         except PermissionError:
#             if attempt < retries - 1:
#                 time.sleep(delay)
#             else:
#                 raise

def getWordsOnlyPDF(path):
    fE = fileExtension(path)
    if (fE != '.pdf'):
        return None
    else:
        text = ''
        try:
            with open(path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page_num in range(3):
                    if page_num < len(reader.pages):    
                        page = reader.pages[page_num]
                        text += page.extract_text()
        except PermissionError as E:
            text = None
        finally:
            return text
    
def getWordsOnlyDOCX(path):
    doc = Document(path)
    fE = fileExtension(path)
    if (fE != '.docx'):
        return None
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    text = ' '.join(text)
    words = text.split()
    words = words[:70]
    words = " ".join(words)
    return words

def getWordsOnlyTXT(path):
    fE = fileExtension(path)
    if (fE not in '.js.java.py.html.css.txt'):
        return None
    with open(path, 'r', encoding = 'utf-8') as file:
        text = file.read()
    return text

def getWordsOnlyPPTX(path):
    ppt = Presentation(path)
    fE = fileExtension(path)
    if (fE != '.pptx'):
        return None
    text = []
    i = 0
    for slide in ppt.slides:
        if i >= 3:
            break
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
        i += 1
    return '\n'.join(text)

def summarizeFileContent(path):
    text_content = ''
    fE = fileExtension(path)
    if (fE == '.pdf'):
        text_content = getWordsOnlyPDF(path)
    elif (fE == '.docx'):
        text_content = getWordsOnlyDOCX(path)
    elif (fE in '.js.java.py.html.css.txt'):
        text_content = getWordsOnlyTXT(path)
    elif (fE == '.pptx'):
        text_content = getWordsOnlyPPTX(path) 
    else:
        return None
    
    features = """
    Please summarise the following text focusing on features that one would consider
    when deciding where to place it in an organised file directory.
    It should be very brief, giving a one to two line summary and some keywords and 
    some guesses as to what the doc could be about. Try to understand what time period/event the doc could signify.
    For example, if the document is titled 'CS 2110 Homework 05', then it could be
    placed into a directory called 'CS 2110' or 'CS 2110/Homeworks' or 'Homeworks/CS 2110'
    What follows are the file name along with some text from the file.
    """
    
    prompt = f"{features}\n\nFilename: {fileName(path)}\nText:\n{text_content}"
    
    response = openai.chat.completions.create(
        model='gpt-4o-mini',
        messages=[
            {'role': 'system', 'content': '''Your role is to summarise a pdf with the given information.'''},
            {'role': 'user', 'content': prompt}
        ],
        max_tokens=60,
        n=1,
        stop=None,
        temperature=0.5,
    )
    summary = response.choices[0].message.content
    return summary
        

def fileContent(path):
    client.set_classifiers(["textrazor_iab_content_taxonomy_3.0"])
    fE = fileExtension(path)
    text_content = ''
    if (fE == '.pdf'):
        text_content = getWordsOnlyPDF(path)
    elif (fE == '.docx'):
        text_content = getWordsOnlyDOCX(path)
    elif (fE in '.js.java.py.html.css.txt'):
        text_content = getWordsOnlyTXT(path)
    elif (fE == '.pptx'):
        text_content = getWordsOnlyPPTX(path) 
    else:
        return None
    text = client.analyze(text_content)
    temp = []
    summary = []

    for c in text.categories():
        if (c.score > 0.4):
            summary.append(c.label.lower())

    for s in text.entities():
        if (s.confidence_score >= 5 and s.relevance_score >= 0.2):
            temp.append(s.matched_text.lower())

    for t in text.topics():
        if (t.score > 0.5):
            temp.append(s.matched_text.lower())

    for p in text.noun_phrases():
        num = 0
        count = 0
        for w in p.words:
            num = num + 1
            if (w.part_of_speech == "NNP" and not w.token.isnumeric() and (w.input_end_offset - w.input_start_offset <= 45) and ((w.input_end_offset - w.input_start_offset >= 2))):
                count = count + 1
        if count == num:
            for w in p.words:
                temp.append(w.token.lower())
    
    for t in temp:
        if (summary.count(t) == 0 and not t.isnumeric()):
            summary.append(t)

    return summary

def webAddress(path):
    return get_web_history(path)

def fileAll(path):
    return fileName(path), fileExtension(path), summarizeFileContent(path), fileContent(path), webAddress(path)